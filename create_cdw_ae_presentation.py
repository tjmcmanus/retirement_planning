#!/usr/bin/env python3
"""
Create CDW Account Executive Presentation for IBM Bob
Why to introduce Bob into $10M+ accounts - Sales enablement focus
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# CDW Color Palette (Primary)
CDW_RED = RGBColor(204, 0, 0)  # #CC0000
CDW_DARK_RED = RGBColor(153, 0, 0)  # #990000
CDW_GRAY = RGBColor(88, 89, 91)  # #58595B
CDW_LIGHT_GRAY = RGBColor(167, 169, 172)  # #A7A9AC

# IBM Color Palette (Partner)
IBM_BLUE = RGBColor(0, 114, 206)  # #0072CE
IBM_DARK_BLUE = RGBColor(0, 67, 206)  # #0043CE

# Accent colors
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(36, 161, 72)
GOLD = RGBColor(255, 193, 7)

def create_presentation():
    """Create the CDW AE presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. Title Slide
    add_title_slide(prs)
    
    # 2. The Opportunity
    add_opportunity_slide(prs)
    
    # 3. Why This Matters to Your Account
    add_account_impact_slide(prs)
    
    # 4. Revenue Opportunity
    add_revenue_opportunity(prs)
    
    # 5. Account Expansion Strategy
    add_expansion_strategy(prs)
    
    # 6. Competitive Positioning
    add_competitive_positioning(prs)
    
    # 7. Customer Pain Points You Can Solve
    add_pain_points_slide(prs)
    
    # 8. Your Value Proposition
    add_value_prop_slide(prs)
    
    # 9. Conversation Starters
    add_conversation_starters(prs)
    
    # 10. Objection Handling
    add_objection_handling(prs)
    
    # 11. Sales Process & Timeline
    add_sales_process(prs)
    
    # 12. Success Metrics
    add_success_metrics(prs)
    
    # 13. Next Steps
    add_next_steps(prs)
    
    return prs

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CDW_RED
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Unlock New Revenue in Your $10M Account"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Why IBM Bob is Your Strategic Account Growth Play"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = GOLD
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Value prop
    value_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.5))
    value_frame = value_box.text_frame
    value_frame.text = "Turn Developer Productivity into\n$500K+ Annual Recurring Revenue"
    for para in value_frame.paragraphs:
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "CDW Account Executive Enablement | IBM Partnership"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = WHITE
    footer_para.alignment = PP_ALIGN.CENTER

def add_opportunity_slide(prs):
    """Add opportunity overview"""
    slide = add_content_slide(prs, "The Opportunity: Why Now?")
    
    # Three key points
    points = [
        ("🎯", "Untapped Wallet Share", "Every developer = $45K annual value\nYour 100-dev account = $4.5M opportunity"),
        ("🚀", "Fast Sales Cycle", "Pilot to production in 60 days\nImmediate ROI proves value"),
        ("🔄", "Recurring Revenue", "Annual subscriptions with high renewal\nExpansion into other business units")
    ]
    
    for idx, (icon, title, desc) in enumerate(points):
        top = Inches(2 + idx * 1.6)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(0.8), Inches(0.8))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_para = icon_frame.paragraphs[0]
        icon_para.font.size = Pt(48)
        icon_para.alignment = PP_ALIGN.CENTER
        
        # Content box
        content_box = slide.shapes.add_textbox(Inches(2), top, Inches(7), Inches(1.2))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
        
        content_frame = content_box.text_frame
        content_frame.text = f"{title}\n{desc}"
        
        para1 = content_frame.paragraphs[0]
        para1.font.size = Pt(22)
        para1.font.bold = True
        para1.font.color.rgb = CDW_RED
        
        if len(content_frame.paragraphs) > 1:
            para2 = content_frame.paragraphs[1]
            para2.font.size = Pt(14)

def add_account_impact_slide(prs):
    """Add account impact slide"""
    slide = add_content_slide(prs, "Why This Matters to Your $10M Account")
    
    # Impact boxes
    impacts = [
        ("Protect Your Base", "Show innovation and value\nPrevent competitive displacement\nDeepen executive relationships"),
        ("Expand Your Footprint", "New budget line (developer tools)\nCross-sell into IT operations\nUpsell existing IBM solutions"),
        ("Increase Your Commission", "$500K+ new ARR opportunity\nHigh-margin software sale\nQuick close = Q4 accelerator")
    ]
    
    for idx, (title, details) in enumerate(impacts):
        left = Inches(0.5 + idx * 3.2)
        
        # Title box
        title_box = slide.shapes.add_textbox(left, Inches(2), Inches(3), Inches(0.6))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = CDW_RED
        
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Details box
        details_box = slide.shapes.add_textbox(left, Inches(2.7), Inches(3), Inches(2.5))
        details_box.fill.solid()
        details_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        details_frame = details_box.text_frame
        details_frame.text = details
        
        for para in details_frame.paragraphs:
            para.font.size = Pt(13)
            para.space_before = Pt(6)
    
    # Bottom banner
    banner = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = IBM_BLUE
    
    banner_frame = banner.text_frame
    banner_frame.text = "Bob positions you as a strategic advisor, not just a vendor"
    banner_para = banner_frame.paragraphs[0]
    banner_para.font.size = Pt(20)
    banner_para.font.bold = True
    banner_para.font.color.rgb = WHITE
    banner_para.alignment = PP_ALIGN.CENTER

def add_revenue_opportunity(prs):
    """Add revenue opportunity slide"""
    slide = add_content_slide(prs, "Your Revenue Opportunity")
    
    # Large number display
    revenue_box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1.5))
    revenue_box.fill.solid()
    revenue_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    
    revenue_frame = revenue_box.text_frame
    revenue_frame.text = "$500K - $2M"
    revenue_para = revenue_frame.paragraphs[0]
    revenue_para.font.size = Pt(72)
    revenue_para.font.bold = True
    revenue_para.font.color.rgb = GREEN
    revenue_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.3), Inches(6), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = "First Year Annual Recurring Revenue"
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(24)
    sub_para.font.bold = True
    sub_para.alignment = PP_ALIGN.CENTER
    
    # Breakdown table
    rows, cols = 5, 3
    left = Inches(1.5)
    top = Inches(4.2)
    width = Inches(7)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Account Size", "Annual Value", "Your Commission*"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CDW_RED
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data
    data = [
        ["50 developers", "$500K", "$25K - $50K"],
        ["100 developers", "$1M", "$50K - $100K"],
        ["200 developers", "$2M", "$100K - $200K"],
        ["+ Services & Training", "+20-30%", "+$10K - $60K"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            if col_idx == 2:  # Commission column
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 205)
                cell.text_frame.paragraphs[0].font.bold = True

def add_expansion_strategy(prs):
    """Add account expansion strategy"""
    slide = add_content_slide(prs, "Your Account Expansion Strategy")
    
    # Strategy steps
    steps = [
        ("Phase 1: Land", "Start with 10-20 developer pilot\n$100K - $200K initial deal\nProve ROI in 30 days"),
        ("Phase 2: Expand", "Roll out to full dev organization\n$500K - $1M expansion\nAdd training & services"),
        ("Phase 3: Multiply", "Expand to other business units\nAdd complementary IBM solutions\n$2M+ total account value")
    ]
    
    for idx, (phase, details) in enumerate(steps):
        top = Inches(2 + idx * 1.5)
        
        # Phase box
        phase_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(2), Inches(0.5))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = IBM_BLUE
        
        phase_frame = phase_box.text_frame
        phase_frame.text = phase
        phase_para = phase_frame.paragraphs[0]
        phase_para.font.size = Pt(18)
        phase_para.font.bold = True
        phase_para.font.color.rgb = WHITE
        phase_para.alignment = PP_ALIGN.CENTER
        
        # Details box
        details_box = slide.shapes.add_textbox(Inches(3), top, Inches(6), Inches(1.2))
        details_frame = details_box.text_frame
        details_frame.text = details
        
        for para in details_frame.paragraphs:
            para.font.size = Pt(14)
            para.space_before = Pt(4)
    
    # Timeline note
    timeline = slide.shapes.add_textbox(Inches(2), Inches(6.2), Inches(6), Inches(0.5))
    timeline_frame = timeline.text_frame
    timeline_frame.text = "Total Timeline: 6-12 months from pilot to full deployment"
    timeline_para = timeline_frame.paragraphs[0]
    timeline_para.font.size = Pt(14)
    timeline_para.font.italic = True
    timeline_para.alignment = PP_ALIGN.CENTER

def add_competitive_positioning(prs):
    """Add competitive positioning"""
    slide = add_content_slide(prs, "Your Competitive Edge")
    
    content = [
        "Why Bob Wins vs. Other AI Tools:",
        "",
        "✓ IBM Partnership = Enterprise Credibility",
        "  Your customer trusts IBM for mission-critical solutions",
        "",
        "✓ CDW Support = Implementation Confidence",
        "  You provide end-to-end service and support",
        "",
        "✓ Proven ROI = Easy Business Case",
        "  $45K per developer is a no-brainer for CFOs",
        "",
        "✓ Fast Time to Value = Quick Wins",
        "  Pilot results in 30 days build momentum",
        "",
        "✓ Exclusive Positioning = No Direct Competition",
        "  GitHub Copilot, Claude, Gemini can't match Bob's capabilities"
    ]
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    text_frame = text_box.text_frame
    
    for i, line in enumerate(content):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        
        if line and not line.startswith(" ") and not line.startswith("✓"):
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = CDW_RED
        elif line.startswith("✓"):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = IBM_BLUE
        else:
            p.font.size = Pt(14)
            p.space_before = Pt(2)

def add_pain_points_slide(prs):
    """Add customer pain points slide"""
    slide = add_content_slide(prs, "Customer Pain Points You Can Solve")
    
    # Pain points with solutions
    pain_points = [
        ("CTO/VP Engineering", "Need to do more with same team\nPressure to accelerate delivery", "Bob delivers 3X productivity"),
        ("CFO", "Rising development costs\nNeed measurable ROI", "Bob shows $45K savings per dev"),
        ("CISO", "Security & compliance concerns\nRisk of AI tools", "Bob is enterprise-grade secure"),
        ("CEO", "Competitive pressure\nTime-to-market critical", "Bob accelerates innovation")
    ]
    
    for idx, (persona, pain, solution) in enumerate(pain_points):
        left = Inches(0.5) if idx < 2 else Inches(5.25)
        top = Inches(2) if idx % 2 == 0 else Inches(4.3)
        
        # Persona box
        persona_box = slide.shapes.add_textbox(left, top, Inches(4.25), Inches(0.4))
        persona_box.fill.solid()
        persona_box.fill.fore_color.rgb = CDW_GRAY
        
        persona_frame = persona_box.text_frame
        persona_frame.text = persona
        persona_para = persona_frame.paragraphs[0]
        persona_para.font.size = Pt(14)
        persona_para.font.bold = True
        persona_para.font.color.rgb = WHITE
        persona_para.alignment = PP_ALIGN.CENTER
        
        # Pain box
        pain_box = slide.shapes.add_textbox(left, top + Inches(0.4), Inches(4.25), Inches(0.8))
        pain_frame = pain_box.text_frame
        pain_frame.text = f"Pain: {pain}"
        for para in pain_frame.paragraphs:
            para.font.size = Pt(11)
        
        # Solution box
        solution_box = slide.shapes.add_textbox(left, top + Inches(1.2), Inches(4.25), Inches(0.5))
        solution_box.fill.solid()
        solution_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
        
        solution_frame = solution_box.text_frame
        solution_frame.text = f"Solution: {solution}"
        solution_para = solution_frame.paragraphs[0]
        solution_para.font.size = Pt(12)
        solution_para.font.bold = True
        solution_para.font.color.rgb = GREEN

def add_value_prop_slide(prs):
    """Add value proposition slide"""
    slide = add_content_slide(prs, "Your Value Proposition to the Customer")
    
    # Value prop statement
    value_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    value_box.fill.solid()
    value_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    value_frame = value_box.text_frame
    value_frame.text = '"We can help you increase developer productivity by 3X while reducing costs by $45,000 per developer annually—with proven results in 30 days."'
    
    for para in value_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.italic = True
        para.alignment = PP_ALIGN.CENTER
        para.font.color.rgb = IBM_DARK_BLUE
    
    # Supporting points
    points = [
        "Backed by IBM's enterprise-grade technology",
        "Supported by CDW's implementation expertise",
        "Proven ROI with Fortune 500 customers",
        "Risk-free pilot program to prove value",
        "Seamless integration with existing tools"
    ]
    
    points_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2.5))
    points_frame = points_box.text_frame
    
    for i, point in enumerate(points):
        if i == 0:
            p = points_frame.paragraphs[0]
        else:
            p = points_frame.add_paragraph()
        
        p.text = f"✓ {point}"
        p.font.size = Pt(16)
        p.space_before = Pt(8)

def add_conversation_starters(prs):
    """Add conversation starters"""
    slide = add_content_slide(prs, "Conversation Starters for Your Next Call")
    
    starters = [
        ('"How satisfied are you with your current developer productivity?"',
         "Opens discussion about pain points"),
        
        ('"What if you could deliver features 70% faster without hiring?"',
         "Quantifies the opportunity"),
        
        ('"We\'re seeing Fortune 500 companies save $45K per developer annually..."',
         "Establishes credibility and ROI"),
        
        ('"Would you be interested in a 30-day pilot with 10 developers?"',
         "Low-risk trial close")
    ]
    
    for idx, (question, note) in enumerate(starters):
        top = Inches(2 + idx * 1.1)
        
        # Question box
        q_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.4), Inches(0.5))
        q_frame = q_box.text_frame
        q_frame.text = question
        q_para = q_frame.paragraphs[0]
        q_para.font.size = Pt(15)
        q_para.font.bold = True
        q_para.font.color.rgb = CDW_RED
        
        # Note box
        note_box = slide.shapes.add_textbox(Inches(1.2), top + Inches(0.5), Inches(7.6), Inches(0.4))
        note_frame = note_box.text_frame
        note_frame.text = f"→ {note}"
        note_para = note_frame.paragraphs[0]
        note_para.font.size = Pt(12)
        note_para.font.italic = True

def add_objection_handling(prs):
    """Add objection handling"""
    slide = add_content_slide(prs, "Handling Common Objections")
    
    objections = [
        ("We already use GitHub Copilot",
         "Bob complements Copilot—it handles multi-file operations and workflows that Copilot can't. Many customers use both."),
        
        ("We need to see ROI first",
         "That's exactly why we offer a 30-day pilot. You'll see measurable results before committing to full deployment."),
        
        ("Our developers are resistant to AI",
         "Bob actually reduces developer frustration by eliminating manual tasks. Early adopters become champions."),
        
        ("Budget is tight this year",
         "Bob pays for itself in 3 months. We can structure payment to align with your budget cycle.")
    ]
    
    for idx, (objection, response) in enumerate(objections):
        top = Inches(2 + idx * 1.1)
        
        # Objection
        obj_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.4), Inches(0.35))
        obj_box.fill.solid()
        obj_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        
        obj_frame = obj_box.text_frame
        obj_frame.text = f"❌ {objection}"
        obj_para = obj_frame.paragraphs[0]
        obj_para.font.size = Pt(13)
        obj_para.font.bold = True
        
        # Response
        resp_box = slide.shapes.add_textbox(Inches(0.8), top + Inches(0.35), Inches(8.4), Inches(0.6))
        resp_frame = resp_box.text_frame
        resp_frame.text = f"✓ {response}"
        resp_para = resp_frame.paragraphs[0]
        resp_para.font.size = Pt(11)
        resp_para.font.color.rgb = GREEN

def add_sales_process(prs):
    """Add sales process slide"""
    slide = add_content_slide(prs, "Your Sales Process & Timeline")
    
    # Process steps
    steps = [
        ("Week 1-2", "Discovery & Qualification", "Identify pain points, stakeholders\nSize the opportunity"),
        ("Week 3-4", "Executive Presentation", "Present business case to CTO/CFO\nPropose pilot program"),
        ("Week 5-6", "Pilot Setup", "10-20 developers, 30 days\nCDW implementation support"),
        ("Week 7-10", "Pilot Results", "Measure ROI, gather testimonials\nExpand business case"),
        ("Week 11-12", "Close & Expand", "Full deployment contract\nServices & training upsell")
    ]
    
    for idx, (timeframe, phase, details) in enumerate(steps):
        top = Inches(1.8 + idx * 0.95)
        
        # Timeline bar
        bar = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.8))
        
        # Alternate colors
        if idx % 2 == 0:
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(240, 248, 255)
        
        frame = bar.text_frame
        frame.text = f"{timeframe}: {phase}\n{details}"
        
        para1 = frame.paragraphs[0]
        para1.font.size = Pt(14)
        para1.font.bold = True
        para1.font.color.rgb = CDW_RED
        
        if len(frame.paragraphs) > 1:
            para2 = frame.paragraphs[1]
            para2.font.size = Pt(11)

def add_success_metrics(prs):
    """Add success metrics slide"""
    slide = add_content_slide(prs, "How to Measure Your Success")
    
    # Metrics boxes
    metrics = [
        ("Deal Size", "$500K - $2M ARR", "Track initial + expansion"),
        ("Sales Cycle", "60-90 days", "Pilot to close"),
        ("Win Rate", "70%+", "After successful pilot"),
        ("Customer Satisfaction", "9/10 NPS", "High renewal rate"),
        ("Account Growth", "3X in Year 2", "Expand to other units"),
        ("Your Commission", "$50K - $200K", "First year earnings")
    ]
    
    for idx, (metric, target, note) in enumerate(metrics):
        row = idx // 3
        col = idx % 3
        left = Inches(0.5 + col * 3.2)
        top = Inches(2 + row * 2)
        
        # Metric box
        box = slide.shapes.add_textbox(left, top, Inches(3), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        frame = box.text_frame
        frame.text = f"{metric}\n{target}\n{note}"
        
        para1 = frame.paragraphs[0]
        para1.font.size = Pt(14)
        para1.font.bold = True
        para1.font.color.rgb = CDW_GRAY
        
        para2 = frame.paragraphs[1]
        para2.font.size = Pt(20)
        para2.font.bold = True
        para2.font.color.rgb = CDW_RED
        
        para3 = frame.paragraphs[2]
        para3.font.size = Pt(11)
        para3.font.italic = True

def add_next_steps(prs):
    """Add next steps slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CDW_RED
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Your Next Steps"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Action items
    actions = [
        "1. Review your $10M account's developer headcount",
        "2. Identify the CTO/VP Engineering as your champion",
        "3. Schedule discovery call to discuss productivity challenges",
        "4. Present Bob as the solution with IBM/CDW backing",
        "5. Propose 30-day pilot with 10-20 developers",
        "6. Close initial deal and plan expansion"
    ]
    
    actions_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3))
    actions_frame = actions_box.text_frame
    
    for i, action in enumerate(actions):
        if i == 0:
            p = actions_frame.paragraphs[0]
        else:
            p = actions_frame.add_paragraph()
        
        p.text = action
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)
    
    # Support
    support_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(7), Inches(0.8))
    support_frame = support_box.text_frame
    support_frame.text = "CDW Sales Engineering & IBM Partnership Team\nare here to support your success"
    
    for para in support_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = GOLD
        para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title):
    """Add a content slide with CDW branding"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = CDW_RED
    header.line.color.rgb = CDW_RED
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Footer bar
    footer = slide.shapes.add_shape(1, Inches(0), Inches(7), Inches(10), Inches(0.5))
    footer.fill.solid()
    footer.fill.fore_color.rgb = CDW_GRAY
    footer.line.color.rgb = CDW_GRAY
    
    # Footer text
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "CDW Account Executive Enablement | IBM Bob Partnership"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = WHITE
    
    return slide

def main():
    """Main function"""
    print("Creating CDW Account Executive Presentation...")
    prs = create_presentation()
    
    filename = "CDW_AE_Bob_Opportunity.pptx"
    prs.save(filename)
    print(f"✅ CDW AE presentation created: {filename}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Focus: Revenue opportunity, account expansion, sales enablement")
    print(f"   Audience: CDW Account Executives with $10M+ accounts")

if __name__ == "__main__":
    main()

# Made with Bob
