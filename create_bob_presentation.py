#!/usr/bin/env python3
"""
Create PowerPoint presentation from Bob Value Proposition markdown
with IBM and CDW color schemes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

# IBM Color Palette
IBM_BLUE = RGBColor(0, 114, 206)  # #0072CE
IBM_DARK_BLUE = RGBColor(0, 67, 206)  # #0043CE
IBM_LIGHT_BLUE = RGBColor(75, 107, 175)  # #4B6BAF
IBM_GRAY = RGBColor(82, 95, 107)  # #525F6B
IBM_DARK_GRAY = RGBColor(22, 22, 22)  # #161616

# CDW Color Palette
CDW_RED = RGBColor(204, 0, 0)  # #CC0000
CDW_DARK_RED = RGBColor(153, 0, 0)  # #990000
CDW_GRAY = RGBColor(88, 89, 91)  # #58595B
CDW_LIGHT_GRAY = RGBColor(167, 169, 172)  # #A7A9AC

# Accent colors
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(36, 161, 72)  # Success green
ORANGE = RGBColor(255, 131, 0)  # Warning orange

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    add_title_slide(prs)
    
    # Executive Summary
    add_executive_summary(prs)
    
    # Architecture Comparison
    add_architecture_comparison(prs)
    
    # Workflow Efficiency
    add_workflow_efficiency(prs)
    
    # Precision & Reliability
    add_precision_reliability(prs)
    
    # Context Management
    add_context_management(prs)
    
    # Real-World Scenarios
    add_real_world_scenarios(prs)
    
    # Advanced Capabilities
    add_advanced_capabilities(prs)
    
    # Safety & Reliability
    add_safety_reliability(prs)
    
    # Integration & Extensibility
    add_integration_extensibility(prs)
    
    # Cost-Effectiveness
    add_cost_effectiveness(prs)
    
    # Use Case Comparison
    add_use_case_comparison(prs)
    
    # Technical Specifications
    add_technical_specs(prs)
    
    # Security & Privacy
    add_security_privacy(prs)
    
    # Future Roadmap
    add_future_roadmap(prs)
    
    # Conclusion
    add_conclusion(prs)
    
    # Summary Comparison
    add_summary_comparison(prs)
    
    # Call to Action
    add_call_to_action(prs)
    
    return prs

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = IBM_DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Bob's Value Proposition"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Why Bob Outperforms Claude, Gemini, and Copilot"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "IBM & CDW Partnership | June 2026"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = CDW_LIGHT_GRAY
    footer_para.alignment = PP_ALIGN.CENTER

def add_executive_summary(prs):
    """Add executive summary slide"""
    slide = add_content_slide(prs, "Executive Summary")
    
    content = [
        "Bob is a specialized AI coding assistant that combines conversational AI with powerful development tools",
        "",
        "Key Differentiators:",
        "• Tool-First Architecture: Direct file system and command execution",
        "• Task-Oriented Workflow: Step-by-step execution with confirmation",
        "• Context-Aware: Deep workspace understanding",
        "• Iterative Development: Built-in feedback loops",
        "• Production-Ready: Designed for real development workflows"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_architecture_comparison(prs):
    """Add architecture comparison slide"""
    slide = add_content_slide(prs, "Architecture Comparison")
    
    # Add comparison table
    rows, cols = 6, 5
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    for i in range(1, 5):
        table.columns[i].width = Inches(1.6)
    
    # Header row
    headers = ["Feature", "Bob", "Claude", "Gemini", "Copilot"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Data rows
    data = [
        ["Direct File Operations", "✅ Native", "❌ Copy/paste", "❌ Copy/paste", "⚠️ Limited"],
        ["Command Execution", "✅ Full CLI", "❌ No", "❌ No", "❌ No"],
        ["Multi-File Editing", "✅ Atomic", "⚠️ Manual", "⚠️ Manual", "⚠️ Single"],
        ["Search & Discovery", "✅ Built-in", "❌ Manual", "❌ Manual", "⚠️ Limited"],
        ["Workspace Context", "✅ Full tree", "⚠️ Manual", "⚠️ Manual", "⚠️ Open files"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Color code Bob column
            if col_idx == 1 and "✅" in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 255, 230)

def add_workflow_efficiency(prs):
    """Add workflow efficiency slide"""
    slide = add_content_slide(prs, "Development Workflow Efficiency")
    
    # Bob's workflow
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(2.5))
    left_frame = left_box.text_frame
    left_frame.text = "With Bob (5 steps):\n\n1. 🔍 search_files\n2. 📖 read_file\n3. ✏️ apply_diff\n4. 🧪 execute_command\n5. ✅ attempt_completion"
    for para in left_frame.paragraphs:
        para.font.size = Pt(14)
    
    # Add green background
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    
    # Others' workflow
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(2.5))
    right_frame = right_box.text_frame
    right_frame.text = "With Claude/Gemini (15+ steps):\n\n1. Ask user to show file\n2. User copies content\n3. Provide suggestions\n4. User copies suggestions\n5. User opens file\n6. User makes edits\n7-15. Repeat..."
    for para in right_frame.paragraphs:
        para.font.size = Pt(14)
    
    # Add red background
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(255, 230, 230)
    
    # Time saved banner
    banner = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.8))
    banner_frame = banner.text_frame
    banner_frame.text = "Time Saved: 70-80% reduction in manual steps"
    banner_para = banner_frame.paragraphs[0]
    banner_para.font.size = Pt(20)
    banner_para.font.bold = True
    banner_para.font.color.rgb = WHITE
    banner_para.alignment = PP_ALIGN.CENTER
    banner.fill.solid()
    banner.fill.fore_color.rgb = CDW_RED

def add_precision_reliability(prs):
    """Add precision & reliability slide"""
    slide = add_content_slide(prs, "Precision & Reliability")
    
    content = [
        "Bob's Targeted Editing Tools:",
        "",
        "apply_diff - Surgical Code Changes",
        "• Exact match verification (no accidental overwrites)",
        "• Multiple changes in one atomic operation",
        "• Line number tracking for precision",
        "• Automatic validation",
        "",
        "Comparison:",
        "• Copilot: Suggests changes, requires manual acceptance",
        "• Claude/Gemini: Full file rewrites (risky for large files)",
        "• Bob: Precise, verified, atomic changes"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_context_management(prs):
    """Add context management slide"""
    slide = add_content_slide(prs, "Context Management")
    
    content = [
        "Bob's Intelligent File Reading:",
        "",
        "Key Features:",
        "• Read up to 5 files simultaneously",
        "• Selective line ranges for large files",
        "• Automatic line numbering",
        "• PDF/DOCX text extraction",
        "",
        "Comparison:",
        "• Claude/Gemini: Limited context window, manual file sharing",
        "• Copilot: Only sees currently open files",
        "• Bob: Proactive context gathering with surgical precision"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_real_world_scenarios(prs):
    """Add real-world scenarios slide"""
    slide = add_content_slide(prs, "Real-World Development Scenarios")
    
    # Scenario 1 table
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Assistant", "Steps", "Time", "Manual Actions"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Data
    data = [
        ["Bob", "4", "2 min", "0"],
        ["Claude", "12+", "15 min", "8+"],
        ["Gemini", "12+", "15 min", "8+"],
        ["Copilot", "8+", "10 min", "5+"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            if row_idx == 1:  # Bob row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 255, 230)
    
    # Add note
    note_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    note_frame = note_box.text_frame
    note_frame.text = "Scenario: Fix a bug affecting 3 files with testing\n\nBob's atomic operations and integrated testing provide 7x faster resolution"
    for para in note_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.italic = True

def add_advanced_capabilities(prs):
    """Add advanced capabilities slide"""
    slide = add_content_slide(prs, "Advanced Capabilities")
    
    content = [
        "Bob's Unique Features:",
        "",
        "1. Code Discovery",
        "   • Instantly map codebase structure",
        "   • Find classes, functions, interfaces",
        "",
        "2. Pattern Search",
        "   • Regex-powered search",
        "   • Context-aware results",
        "",
        "3. Task Management",
        "   • Built-in task tracking",
        "   • Progress visualization",
        "",
        "None of these exist in Claude, Gemini, or Copilot"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_safety_reliability(prs):
    """Add safety & reliability slide"""
    slide = add_content_slide(prs, "Safety & Reliability")
    
    # Safety features table
    rows, cols = 6, 3
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(3.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Feature", "Bob", "Others"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Data
    data = [
        ["Exact Match Verification", "✅ Required", "❌ No"],
        ["Atomic Operations", "✅ All or nothing", "⚠️ Partial"],
        ["User Confirmation", "✅ Every step", "⚠️ Optional"],
        ["Rollback Support", "✅ Via version control", "⚠️ Manual"],
        ["File Locking", "✅ .bobignore", "❌ No"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(11)

def add_integration_extensibility(prs):
    """Add integration & extensibility slide"""
    slide = add_content_slide(prs, "Integration & Extensibility")
    
    content = [
        "Bob's Ecosystem:",
        "",
        "Core Engine:",
        "• File Operations • Command Execution",
        "• Search & Discovery • Code Analysis",
        "",
        "VS Code Integration:",
        "• Workspace Awareness • Terminal Access",
        "• Git Integration • Extension Ecosystem",
        "",
        "Development Tools:",
        "• npm/pip/cargo/etc. • Testing Frameworks",
        "• Linters & Formatters • Build Systems",
        "",
        "Bob provides full development environment access"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_cost_effectiveness(prs):
    """Add cost-effectiveness slide"""
    slide = add_content_slide(prs, "Cost-Effectiveness & ROI")
    
    # ROI box
    roi_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(2))
    roi_frame = roi_box.text_frame
    roi_frame.text = "ROI Calculation:\n\nDeveloper hourly rate: $75/hour\nTime saved per month: 50 hours\nMonthly value: $3,750\n\nAnnual value: $45,000 per developer"
    for para in roi_frame.paragraphs:
        para.font.size = Pt(16)
        if "$45,000" in para.text:
            para.font.bold = True
            para.font.size = Pt(20)
            para.font.color.rgb = CDW_RED
    
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # Time savings
    savings_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.5))
    savings_frame = savings_box.text_frame
    savings_frame.text = "Time Savings Per Task:\n• Simple edits: 50% faster\n• Multi-file refactoring: 70% faster\n• Bug investigation: 60% faster\n• Feature implementation: 65% faster"
    for para in savings_frame.paragraphs:
        para.font.size = Pt(14)

def add_use_case_comparison(prs):
    """Add use case comparison slide"""
    slide = add_content_slide(prs, "When to Use Each Tool")
    
    # Create 4 quadrants
    quadrants = [
        ("Bob - Best For:", [
            "✅ Multi-file refactoring",
            "✅ Bug fixes requiring investigation",
            "✅ Feature implementation",
            "✅ Code migration",
            "✅ Testing and validation"
        ], Inches(0.5), Inches(2)),
        ("Claude - Best For:", [
            "✅ Brainstorming and design",
            "✅ Code review and explanation",
            "✅ Algorithm design",
            "⚠️ Requires manual file operations"
        ], Inches(5.25), Inches(2)),
        ("Gemini - Best For:", [
            "✅ Research and learning",
            "✅ Code explanation",
            "✅ General programming questions",
            "⚠️ Limited coding workflow support"
        ], Inches(0.5), Inches(4.5)),
        ("Copilot - Best For:", [
            "✅ Inline code completion",
            "✅ Single-file edits",
            "⚠️ No multi-file coordination",
            "⚠️ No command execution"
        ], Inches(5.25), Inches(4.5))
    ]
    
    for title, items, left, top in quadrants:
        box = slide.shapes.add_textbox(left, top, Inches(4.25), Inches(2))
        frame = box.text_frame
        frame.text = title + "\n" + "\n".join(items)
        for para in frame.paragraphs:
            para.font.size = Pt(11)
            if title.startswith("Bob"):
                para.font.bold = True if para.text == title else False

def add_technical_specs(prs):
    """Add technical specifications slide"""
    slide = add_content_slide(prs, "Technical Specifications")
    
    content = [
        "Core Components:",
        "• LLM Engine: Advanced language understanding",
        "• Tool System: 12+ specialized development tools",
        "• File System Interface: Direct OS-level access",
        "• Command Executor: Full shell integration",
        "",
        "Performance Metrics:",
        "• File operations: <100ms",
        "• Search operations: <500ms",
        "• Multi-file edits: <2s",
        "",
        "Supported Languages:",
        "Python, JavaScript, TypeScript, Java, C++, C#, Go, Rust,",
        "Ruby, PHP, Swift, Kotlin, and 50+ more"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_security_privacy(prs):
    """Add security & privacy slide"""
    slide = add_content_slide(prs, "Security & Privacy")
    
    content = [
        "Bob's Security Model:",
        "",
        "Data Protection:",
        "✅ Local execution (no code sent to cloud)",
        "✅ File access controls (.bobignore)",
        "✅ User confirmation required",
        "✅ Audit trail of all operations",
        "✅ No persistent storage of code",
        "",
        "Comparison:",
        "• Copilot: Sends code to GitHub servers",
        "• Claude/Gemini: Processes in cloud",
        "• Bob: Hybrid model with local tool execution"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_future_roadmap(prs):
    """Add future roadmap slide"""
    slide = add_content_slide(prs, "Future Roadmap")
    
    content = [
        "Q3 2026:",
        "• MCP (Model Context Protocol) integration",
        "• Custom tool creation",
        "• Advanced analytics dashboard",
        "• Git workflow automation",
        "",
        "Q4 2026:",
        "• Multi-repository support",
        "• Automated testing frameworks",
        "• Documentation generation",
        "• CI/CD integration",
        "",
        "2027:",
        "• Team collaboration features",
        "• Performance profiling",
        "• Enhanced security scanning"
    ]
    
    add_bullet_content(slide, content, Inches(1), Inches(2), Inches(8), Inches(4.5))

def add_conclusion(prs):
    """Add conclusion slide"""
    slide = add_content_slide(prs, "Why Bob is the Superior Choice")
    
    # Three columns
    columns = [
        ("For Individual Developers:", [
            "⚡ 60-70% faster workflow",
            "🎯 Reduced context switching",
            "✅ Higher code quality",
            "🧘 Less cognitive load"
        ]),
        ("For Teams:", [
            "📈 Increased productivity",
            "🤝 Consistent code patterns",
            "📚 Better knowledge sharing",
            "💰 Significant cost savings"
        ]),
        ("For Organizations:", [
            "💵 $45K+ ROI per developer",
            "🚀 Faster time to market",
            "🛡️ Reduced bug rates",
            "📊 Measurable productivity gains"
        ])
    ]
    
    for idx, (title, items) in enumerate(columns):
        left = Inches(0.5 + idx * 3.2)
        box = slide.shapes.add_textbox(left, Inches(2), Inches(3), Inches(4))
        frame = box.text_frame
        frame.text = title + "\n\n" + "\n".join(items)
        for para in frame.paragraphs:
            para.font.size = Pt(12)
            if para.text == title:
                para.font.bold = True
                para.font.size = Pt(14)

def add_summary_comparison(prs):
    """Add summary comparison slide"""
    slide = add_content_slide(prs, "Summary: The Bob Advantage")
    
    # Comparison table
    rows, cols = 11, 5
    left = Inches(0.3)
    top = Inches(1.8)
    width = Inches(9.4)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Capability", "Bob", "Claude", "Gemini", "Copilot"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Data
    data = [
        ["Direct File Editing", "⭐⭐⭐⭐⭐", "⭐", "⭐", "⭐⭐⭐"],
        ["Multi-File Operations", "⭐⭐⭐⭐⭐", "⭐⭐", "⭐⭐", "⭐⭐"],
        ["Command Execution", "⭐⭐⭐⭐⭐", "⭐", "⭐", "⭐"],
        ["Code Search", "⭐⭐⭐⭐⭐", "⭐⭐", "⭐⭐", "⭐⭐⭐"],
        ["Context Awareness", "⭐⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐"],
        ["Workflow Integration", "⭐⭐⭐⭐⭐", "⭐⭐", "⭐⭐", "⭐⭐⭐⭐"],
        ["Safety & Reliability", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐"],
        ["Conversation Quality", "⭐⭐⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐"],
        ["Code Completion", "⭐⭐⭐", "⭐⭐", "⭐⭐", "⭐⭐⭐⭐⭐"],
        ["Overall Development", "⭐⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Highlight Bob's 5-star ratings
            if col_idx == 1 and "⭐⭐⭐⭐⭐" in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 250, 205)

def add_call_to_action(prs):
    """Add call to action slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background gradient effect
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = IBM_DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Try Bob Today"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Getting started
    steps_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    steps_frame = steps_box.text_frame
    steps_frame.text = "Getting Started:\n\n1. Install Bob in VS Code\n2. Open your project\n3. Give Bob a task\n4. Watch the magic happen"
    for para in steps_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Resources:\n📖 Documentation | 💬 Community | 🎓 Tutorials | 🐛 Support"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = CDW_LIGHT_GRAY
        para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title):
    """Add a content slide with IBM/CDW branding"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add header bar
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = IBM_BLUE
    header.line.color.rgb = IBM_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Add footer bar
    footer = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(7),
        Inches(10), Inches(0.5)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = CDW_GRAY
    footer.line.color.rgb = CDW_GRAY
    
    # Add footer text
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Bob Value Proposition | IBM & CDW"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = WHITE
    
    return slide

def add_bullet_content(slide, content_list, left, top, width, height):
    """Add bullet point content to a slide"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.size = Pt(14)
        
        # Adjust formatting based on content
        if line and not line.startswith(" ") and not line.startswith("•"):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = IBM_DARK_BLUE
        elif line.startswith("•"):
            p.level = 0
        elif line.startswith("  "):
            p.level = 1
            p.font.size = Pt(12)

def main():
    """Main function"""
    print("Creating Bob Value Proposition presentation...")
    prs = create_presentation()
    
    filename = "Bob_Value_Proposition.pptx"
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Color scheme: IBM Blue & CDW Red")

if __name__ == "__main__":
    main()

# Made with Bob
